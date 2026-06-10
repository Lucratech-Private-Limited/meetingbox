"""
Gmail Service -- send emails using stored OAuth2 tokens.
"""

import base64
import logging
import os
import re
import time
from datetime import datetime, timedelta, timezone
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.utils import getaddresses

from googleapiclient.discovery import build

logger = logging.getLogger(__name__)

# Gmail search: Primary + explicit calendar/ICS branch (feeds the allowlist with candidates).
_SUBJECT_NOISE_IN_QUERY = (
    "unsubscribe OR newsletter OR \"security alert\" OR \"verification code\" OR "
    "\"password reset\" OR \"two-factor\" OR \"2-step verification\" OR \"payment received\" OR "
    "\"your receipt\" OR \"subscription renew\" OR \"subscription confirmation\" OR invoice OR "
    "\"order confirmation\" OR \"billing statement\" OR \"weekly digest\""
)

_DEFAULT_GMAIL_LIST_Q = (
    "("
    "(in:inbox category:primary -category:promotions -category:social -category:forums "
    f"-subject:({_SUBJECT_NOISE_IN_QUERY}))"
    " OR "
    "(in:inbox (filename:ics OR from:calendar-notification@google.com OR from:reminders@google.com OR "
    "subject:(invitation OR \"Updated invitation\" OR \"Event reminder\" OR "
    "\"Meeting reminder\" OR \"Reminder:\")))"
    ")"
)

_FROM_EMAIL_RE = re.compile(r"[\w.+-]+@[\w.-]+\.[a-zA-Z]{2,}")

# Gmail category tabs — detected from API message labelIds (no domain blocklists).
_BULK_TAB_LABELS = frozenset({
    "CATEGORY_PROMOTIONS",
    "CATEGORY_SOCIAL",
    "CATEGORY_FORUMS",
    "CATEGORY_UPDATES",
})

# Optional extra domains (admin/env only); default is empty — see MEETINGBOX_GMAIL_BLOCK_DOMAINS.
_DEFAULT_BLOCKED_DOMAINS: frozenset[str] = frozenset()

_CALENDAR_MEETING_REMINDER = re.compile(
    r"invitation|accepted:|declined:|tentative:|updated invitation|event reminder|meeting reminder|"
    r"^\s*reminder:|\.ics\b|\bcalendar invite\b",
    re.I,
)
_TASKISH_SUBJECT = re.compile(
    r"\btodo\b|\btask(s)? due\b|\btask reminder\b|assigned to you|action item",
    re.I,
)
_NEWSLETTER_SNIPPET = re.compile(
    r"view (this )?email in (your )?browser|unsubscribe|manage (your )?preferences|email preferences|"
    r"you(’|')re receiving|mailing list",
    re.I,
)
_NOISE_SUBJECT = re.compile(
    r"security alert|verification code|password (was )?reset|two[- ]?factor|otp\b|"
    r"payment (received|successful|due)|your receipt|billing statement|subscription (renew|confirm)|"
    r"order confirmation",
    re.I,
)
_TRANSACTIONAL_FROM = re.compile(
    r"noreply|no-reply|donotreply|mailer-daemon|notifications?@|"
    r"billing@|payments?@|subscriptions?@|"
    r"@email\.(stripe|paypal)|@(stripe|paypal)\.com",
    re.I,
)

# Local-parts that are only used for automated mail (exact match).
_BLOCKED_LOCAL_PARTS_EXACT = frozenset({
    "noreply",
    "no-reply",
    "donotreply",
    "donot-reply",
    "mailer-daemon",
    "postmaster",
    "newsletter",
    "notifications",
    "notification",
    "notify",
    "alerts",
    "billing",
    "payments",
    "invoices",
    "invoice",
    "orders",
    "subscribe",
    "marketing",
    "promo",
    "team",
    "mail",
})

# Compound local-parts (e.g. antigravity-noreply@): pattern-based, not domain blocklists.
_LOCAL_AUTOMATED_RE = re.compile(
    r"^no-?reply([+._]|$)|"
    r"^donotreply|^donot-reply|^mailer-daemon|^postmaster|"
    r"^bounce|^bounces|"
    r"[-.+_]no-?reply$|"
    r"(^|[-.+_])noreply($|[-.+_@])|"
    r"^notifications?|^notify|^alerts?|^digest$|^news$|^mailshot|"
    r"^feedback|^support\+|^mailing|^promo|^subscribe",
    re.I,
)


def _blocked_domains() -> frozenset[str]:
    """Optional admin domain blocklist only (MEETINGBOX_GMAIL_BLOCK_DOMAINS). Empty by default."""
    raw = (os.getenv("MEETINGBOX_GMAIL_BLOCK_DOMAINS") or "").strip()
    if not raw:
        return _DEFAULT_BLOCKED_DOMAINS
    extra = {p.strip().lower() for p in raw.split(",") if p.strip()}
    return _DEFAULT_BLOCKED_DOMAINS | extra


def _headers_lc_from_meta(headers: list | None) -> dict[str, str]:
    out: dict[str, str] = {}
    for h in headers or []:
        if not isinstance(h, dict):
            continue
        name = (h.get("name") or "").strip().lower()
        if name:
            out[name] = (h.get("value") or "").strip()
    return out


def _is_bulk_or_automated_clutter(
    *,
    from_hdr: str,
    subject: str,
    snippet: str,
    headers_lc: dict[str, str] | None,
    label_ids: list[str] | None,
) -> bool:
    """
    Marketing / transactional / list mail using Gmail labels and RFC bulk signals,
    plus automated From local-parts (no static domain denylist).
    """
    lids = {str(x) for x in (label_ids or [])}
    if lids & _BULK_TAB_LABELS:
        return True

    h = headers_lc or {}
    if h.get("list-unsubscribe") or h.get("list-unsubscribe-post"):
        return True
    if h.get("list-id"):
        return True

    prec = (h.get("precedence") or "").lower().strip()
    if prec in ("bulk", "list", "junk"):
        return True

    auto_sub = (h.get("auto-submitted") or "").lower().strip()
    if auto_sub and auto_sub != "no":
        if "auto-generated" in auto_sub or "auto-replied" in auto_sub:
            return True

    if h.get("feedback-id"):
        return True

    for key in h:
        kl = key.lower()
        if kl.startswith("x-mailgun") or kl.startswith("x-ses-") or kl == "x-campaignid":
            return True
        if kl in ("x-mailchimp", "x-mc-user", "x-sendgrid", "x-sendinblue"):
            return True

    addr = _extract_address(from_hdr)
    if addr:
        local, _, domain = addr.partition("@")
        if local in _BLOCKED_LOCAL_PARTS_EXACT:
            return True
        if _local_part_looks_automated(local):
            return True
        blocked = _blocked_domains()
        if _domain_suffix_blocked(domain, blocked):
            return True

    if _TRANSACTIONAL_FROM.search(from_hdr or ""):
        return True

    if _NEWSLETTER_SNIPPET.search(snippet or "") or _NEWSLETTER_SNIPPET.search(subject or ""):
        return True

    if _NOISE_SUBJECT.search(subject or ""):
        return True

    return False


def _local_part_looks_automated(local: str) -> bool:
    if not local:
        return False
    lo = local.lower()
    if lo in _BLOCKED_LOCAL_PARTS_EXACT:
        return True
    return bool(_LOCAL_AUTOMATED_RE.search(local))


def _strict_allowlist_enabled() -> bool:
    return (os.getenv("MEETINGBOX_GMAIL_STRICT_INBOX", "") or "").strip() != "0"


def _extract_address(from_hdr: str) -> str | None:
    m = _FROM_EMAIL_RE.search(from_hdr or "")
    return m.group(0).lower() if m else None


def _domain_suffix_blocked(domain: str, blocked: frozenset[str]) -> bool:
    d = domain.lower()
    for b in blocked:
        if d == b or d.endswith("." + b):
            return True
    return False


def _is_calendar_or_system_reminder(from_hdr: str, subject: str, snippet: str) -> bool:
    f = (from_hdr or "").lower()
    s = (subject or "").lower()
    sn = (snippet or "").lower()

    if "calendar-notification@google.com" in f or "reminders@google.com" in f:
        return True

    if "outlook.com" in f or "microsoft.com" in f:
        if _CALENDAR_MEETING_REMINDER.search(s):
            return True

    # Strong calendar / meeting signals (avoid bare "invitation" = LinkedIn etc.).
    if re.search(
        r"accepted:|declined:|tentative:|updated invitation|^invitation:|\.ics\b|"
        r"event reminder|meeting reminder|^\s*reminder:|calendar invite",
        s,
        re.I,
    ):
        return True

    if _TASKISH_SUBJECT.search(s) or _TASKISH_SUBJECT.search(sn):
        return True

    return False


def should_show_in_personal_inbox(
    from_hdr: str,
    subject: str,
    snippet: str,
    *,
    headers_lc: dict[str, str] | None = None,
    label_ids: list[str] | None = None,
) -> bool:
    """
    Human + meeting/task mail: calendar exceptions first, then hide list/bulk/automated
    using Gmail labelIds, RFC list headers, and From local-part heuristics (no domain denylist by default).
    """
    subj = subject or ""
    snip = snippet or ""
    frm = from_hdr or ""

    if _is_calendar_or_system_reminder(frm, subj, snip):
        return True

    if _is_bulk_or_automated_clutter(
        from_hdr=frm,
        subject=subj,
        snippet=snip,
        headers_lc=headers_lc,
        label_ids=label_ids,
    ):
        return False

    if not _extract_address(frm):
        return False

    return True


# Legacy soft-hide when strict mode is off (same clutter rules, or legacy noise+transactional only when clutter checks are inconclusive).
def _postfilter_disabled() -> bool:
    return (os.getenv("MEETINGBOX_GMAIL_POSTFILTER", "") or "").strip() == "0"


def _should_hide_list_row_soft(
    subject: str,
    from_hdr: str,
    snippet: str = "",
    *,
    headers_lc: dict[str, str] | None = None,
    label_ids: list[str] | None = None,
) -> bool:
    if _postfilter_disabled():
        return False
    subj = subject or ""
    frm = from_hdr or ""
    snip = snippet or ""
    if _is_calendar_or_system_reminder(frm, subj, snip):
        return False
    if _is_bulk_or_automated_clutter(
        from_hdr=frm,
        subject=subj,
        snippet=snip,
        headers_lc=headers_lc,
        label_ids=label_ids,
    ):
        return True
    # Legacy narrow path
    if _CALENDAR_MEETING_REMINDER.search(subj):
        return False
    if not _NOISE_SUBJECT.search(subj):
        return False
    if not _TRANSACTIONAL_FROM.search(frm):
        return False
    return True


def default_gmail_list_query() -> str:
    """Override entire default search string via MEETINGBOX_GMAIL_LIST_DEFAULT_Q."""
    raw = (os.getenv("MEETINGBOX_GMAIL_LIST_DEFAULT_Q") or "").strip()
    return raw if raw else _DEFAULT_GMAIL_LIST_Q


def _merge_disabled() -> bool:
    return (os.getenv("MEETINGBOX_GMAIL_LIST_MERGE", "") or "").strip() == "0"


def _should_merge_inbox_scope(user_q: str) -> bool:
    if _merge_disabled():
        return False
    lower = user_q.lower()
    skip_tokens = (
        "in:spam",
        "in:trash",
        "in:all",
        "in:anywhere",
        "in:snoozed",
        "in:important",
        "in:sent",
        "in:drafts",
    )
    return not any(tok in lower for tok in skip_tokens)


def compose_list_query(user_q: str) -> str:
    q = (user_q or "").strip()
    base = default_gmail_list_query()
    if not q:
        return base
    if not _should_merge_inbox_scope(q):
        return q
    if base == q:
        return q
    return f"({q}) ({base})"


def gmail_search_after_clause(days: int) -> str:
    """Gmail ``q`` fragment: messages on or after this calendar day (UTC)."""
    cutoff = (datetime.now(timezone.utc) - timedelta(days=max(1, int(days)))).date()
    return f"after:{cutoff.year}/{cutoff.month:02d}/{cutoff.day:02d}"


def merge_days_into_gmail_query(user_q: str, days: int | None) -> str:
    """When *days* is set, prefix the query with Gmail ``after:`` (bounded listing)."""
    if days is None:
        return user_q or ""
    after_part = gmail_search_after_clause(days)
    u = (user_q or "").strip()
    return f"{after_part} {u}".strip() if u else after_part


def list_recent_messages(
    credentials,
    max_results: int = 10,
    q: str = "",
    days: int | None = None,
) -> list[dict]:
    """
    List recent message metadata (From, Subject, Date, snippet, threadId).
    Requires gmail.readonly scope.

    By default applies MEETINGBOX_GMAIL_STRICT_INBOX allowlist so UIs only see
    individual + meeting/task/reminder mail. Set MEETINGBOX_GMAIL_STRICT_INBOX=0 for legacy behavior.

    If *days* is set, the Gmail search is scoped with ``after:YYYY/MM/DD`` (UTC) so the API
    returns a bounded window (e.g. dashboard inbox lists).
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    max_results = max(1, min(int(max_results), 50))
    q_merged = merge_days_into_gmail_query(q, days)
    q_use = compose_list_query(q_merged)
    strict = _strict_allowlist_enabled()

    logger.debug(
        "Gmail messages.list strict=%s q=%s",
        strict,
        q_use[:500] + ("…" if len(q_use) > 500 else ""),
    )

    out: list[dict] = []
    seen_ids: set[str] = set()
    page_token: str | None = None
    rounds = 0
    max_rounds = 15
    max_inspect = min(500, max(80, max_results * 25))
    inspected = 0

    list_batch = 40

    _META_HEADERS = [
        "From",
        "Subject",
        "Date",
        "List-Unsubscribe",
        "List-Unsubscribe-Post",
        "List-Id",
        "Precedence",
        "Auto-Submitted",
        "Feedback-ID",
    ]

    while len(out) < max_results and rounds < max_rounds and inspected < max_inspect:
        kwargs: dict = {"userId": "me", "maxResults": list_batch, "q": q_use}
        if page_token:
            kwargs["pageToken"] = page_token
        results = service.users().messages().list(**kwargs).execute()
        msg_refs = results.get("messages", [])
        page_token = results.get("nextPageToken")
        rounds += 1

        if not msg_refs:
            break

        # Collect this page's ids in list order (newest-first), de-duped and
        # capped by max_inspect — identical selection to the previous per-message
        # loop, just gathered up front so the metadata gets can be batched.
        page_ids: list[str] = []
        for ref in msg_refs:
            mid = ref.get("id")
            if not mid or mid in seen_ids:
                continue
            seen_ids.add(mid)
            inspected += 1
            if inspected > max_inspect:
                break
            page_ids.append(mid)

        # Fetch all metadata for this page in ONE batched HTTP request instead of
        # one sequential round-trip per message (the prior N+1 that dominated
        # latency). Results may arrive out of order, so key them by message id and
        # then walk page_ids in original order to preserve identical output.
        fetched: dict[str, dict] = {}

        def _collect(request_id, response, exception, _store=fetched):
            if exception is None and response is not None:
                _store[request_id] = response

        # Fetch this page's metadata in batched rounds, retrying only the ids
        # that did not come back (transient 429/5xx under load). Without the
        # retry a dropped message would be silently skipped and the early-stop
        # below would substitute an OLDER message, yielding an incomplete and
        # unstable inbox view. Retrying keeps the result the true newest-N set.
        pending = list(page_ids)
        for _attempt in range(3):
            if not pending:
                break
            if _attempt:
                time.sleep(0.4 * _attempt)
            batch = service.new_batch_http_request(callback=_collect)
            for mid in pending:
                batch.add(
                    service.users().messages().get(
                        userId="me",
                        id=mid,
                        format="metadata",
                        metadataHeaders=_META_HEADERS,
                    ),
                    request_id=mid,
                )
            batch.execute()
            pending = [mid for mid in pending if mid not in fetched]
        if pending:
            logger.warning(
                "gmail list: %d/%d message metadata fetch(es) failed after retries",
                len(pending),
                len(page_ids),
            )

        for mid in page_ids:
            if len(out) >= max_results:
                break
            m = fetched.get(mid)
            if not m:
                continue
            raw_headers = m.get("payload", {}).get("headers", [])
            headers_lc = _headers_lc_from_meta(raw_headers)
            label_ids = list(m.get("labelIds") or [])
            subj = headers_lc.get("subject", "")
            frm = headers_lc.get("from", "")
            snippet = m.get("snippet", "") or ""

            if strict:
                if not should_show_in_personal_inbox(
                    frm,
                    subj,
                    snippet,
                    headers_lc=headers_lc,
                    label_ids=label_ids,
                ):
                    continue
            elif _should_hide_list_row_soft(
                subj,
                frm,
                snippet,
                headers_lc=headers_lc,
                label_ids=label_ids,
            ):
                continue

            out.append({
                "id": mid,
                "threadId": m.get("threadId"),
                "snippet": snippet,
                "from": frm,
                "subject": subj,
                "date": headers_lc.get("date", ""),
                "is_read": "UNREAD" not in label_ids,
            })

        if not page_token:
            break

    return out


def get_self_email(credentials) -> str:
    """Return the authenticated user's own Gmail address (lower-cased)."""
    try:
        service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
        prof = service.users().getProfile(userId="me").execute()
        return str(prof.get("emailAddress") or "").strip().lower()
    except Exception as exc:
        logger.debug("get_self_email failed: %s", exc)
        return ""


def harvest_contact_addresses(
    credentials,
    query: str = "",
    max_messages: int = 200,
) -> list[str]:
    """Collect raw From / To / Cc header strings across the user's mailbox.

    Unlike ``list_recent_messages`` this fetches the recipient headers (To/Cc)
    and applies NO personal-inbox allowlist — the goal is to learn EVERY person
    the user has corresponded with (sent or received), so the contacts book can
    resolve a name regardless of direction.

    Returns a flat list of address-field strings (e.g. 'Jane Roe <jane@x.com>')
    suitable for ``email.utils.getaddresses`` / contacts parsing.
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    max_messages = max(1, min(int(max_messages), 1000))
    out: list[str] = []
    seen: set[str] = set()
    page_token: str | None = None
    fetched = 0

    while fetched < max_messages:
        kwargs: dict = {
            "userId": "me",
            "maxResults": min(100, max_messages - fetched),
            "q": query or "",
        }
        if page_token:
            kwargs["pageToken"] = page_token
        try:
            resp = service.users().messages().list(**kwargs).execute()
        except Exception as exc:
            logger.debug("harvest list failed: %s", exc)
            break
        refs = resp.get("messages", []) or []
        page_token = resp.get("nextPageToken")
        if not refs:
            break
        for ref in refs:
            mid = ref.get("id")
            if not mid or mid in seen:
                continue
            seen.add(mid)
            fetched += 1
            try:
                m = (
                    service.users()
                    .messages()
                    .get(
                        userId="me",
                        id=mid,
                        format="metadata",
                        metadataHeaders=["From", "To", "Cc"],
                    )
                    .execute()
                )
            except Exception:
                continue
            for h in m.get("payload", {}).get("headers", []) or []:
                if (h.get("name") or "").lower() in ("from", "to", "cc"):
                    val = h.get("value") or ""
                    if val.strip():
                        out.append(val)
        if not page_token:
            break

    return out


def send_email(
    credentials,
    to: str,
    subject: str,
    body: str,
    html_body: str | None = None,
    cc: str | None = None,
    bcc: str | None = None,
    thread_id: str | None = None,
) -> dict:
    """
    Send an email via Gmail API.

    Args:
        credentials: google.oauth2.credentials.Credentials with gmail.send scope
        to: recipient address
        subject: subject line
        body: plain text body
        html_body: optional HTML body
        cc: optional CC
        bcc: optional BCC

    Returns:
        Google API response dict with 'id', 'threadId', 'labelIds'
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)

    if html_body:
        message = MIMEMultipart("alternative")
        message.attach(MIMEText(body, "plain"))
        message.attach(MIMEText(html_body, "html"))
    else:
        message = MIMEText(body, "plain")

    message["to"] = to
    message["subject"] = subject
    if cc:
        message["cc"] = cc
    if bcc:
        message["bcc"] = bcc

    raw = base64.urlsafe_b64encode(message.as_bytes()).decode("utf-8")

    send_body: dict = {"raw": raw}
    if thread_id and str(thread_id).strip():
        send_body["threadId"] = str(thread_id).strip()

    result = (
        service.users()
        .messages()
        .send(userId="me", body=send_body)
        .execute()
    )

    logger.info("Email sent: id=%s to=%s subject=%s", result.get("id"), to, subject)
    return result


def create_draft(
    credentials,
    to: str = "",
    subject: str = "",
    body: str = "",
    cc: str | None = None,
) -> dict:
    """
    Create a Gmail draft (does not send). Requires gmail.compose scope.
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    message = MIMEText(body, "plain")
    message["to"] = to
    message["subject"] = subject
    if cc:
        message["cc"] = cc
    raw = base64.urlsafe_b64encode(message.as_bytes()).decode("utf-8")
    draft_body = {"message": {"raw": raw}}
    result = (
        service.users()
        .drafts()
        .create(userId="me", body=draft_body)
        .execute()
    )
    logger.info("Gmail draft created: id=%s subject=%s", result.get("id"), subject)
    return result


def send_draft(credentials, draft_id: str) -> dict:
    """
    Send an existing Gmail draft by draft_id using the drafts.send API.
    Gmail automatically moves it from Drafts to Sent and removes it from Drafts.

    Also reads the draft's recipient headers (To/Cc/Bcc) BEFORE sending and
    returns them as ``recipients`` so the caller can remember confirmed
    contacts — addresses are learned only on a real send, never from a draft
    that may still hold a mis-heard address.
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)

    recipients: list[str] = []
    try:
        draft = (
            service.users()
            .drafts()
            .get(userId="me", id=draft_id, format="metadata")
            .execute()
        )
        env = _parse_gmail_message_envelope(draft.get("message", {}))
        recipients = [v for v in (env.get("to"), env.get("cc"), env.get("bcc")) if v]
    except Exception as exc:
        logger.debug("send_draft: could not read recipients for %s: %s", draft_id, exc)

    result = (
        service.users()
        .drafts()
        .send(userId="me", body={"id": draft_id})
        .execute()
    )
    logger.info(
        "Gmail draft sent: draft_id=%s message_id=%s thread_id=%s",
        draft_id,
        result.get("id"),
        result.get("threadId"),
    )
    return {
        "id": result.get("id"),
        "thread_id": result.get("threadId"),
        "draft_id": draft_id,
        "status": "sent",
        "recipients": recipients,
    }


def get_message_full(credentials, message_id: str) -> dict:
    """
    Fetch a single message with full payload (headers + body text).
    Returns dict with id, from, subject, date, body, snippet, is_read.
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    m = service.users().messages().get(
        userId="me",
        id=message_id,
        format="full",
    ).execute()

    headers = {h["name"]: h["value"] for h in m.get("payload", {}).get("headers", [])}
    label_ids = m.get("labelIds", [])
    is_read = "UNREAD" not in label_ids

    # Extract plain-text body
    body = _extract_body(m.get("payload", {}))

    return {
        "id": m["id"],
        "threadId": m.get("threadId"),
        "from": headers.get("From", ""),
        "to": headers.get("To", ""),
        "subject": headers.get("Subject", ""),
        "date": headers.get("Date", ""),
        "snippet": m.get("snippet", ""),
        "body": body or m.get("snippet", ""),
        "is_read": is_read,
        "label_ids": label_ids,
    }


def _extract_body(payload: dict) -> str:
    """Recursively extract the first text/plain part from a Gmail payload."""
    import base64

    mime = payload.get("mimeType", "")
    if mime == "text/plain":
        data = payload.get("body", {}).get("data", "")
        if data:
            try:
                return base64.urlsafe_b64decode(data + "==").decode("utf-8", errors="replace")
            except Exception:
                return ""

    for part in payload.get("parts", []):
        result = _extract_body(part)
        if result:
            return result
    return ""


_ATTACHMENT_MAX_TEXT_BYTES = 50_000  # 50 KB of extracted text per attachment
_ATTACHMENT_MAX_COUNT = 5            # max attachments to process per message


def _extract_attachments_metadata(payload: dict) -> list[dict]:
    """
    Recursively walk a Gmail message payload and collect every attachment part.
    A part is an attachment when it carries a non-empty filename.
    Returns list of {attachment_id, data_inline, filename, mime_type, size}.
    data_inline is the base64-encoded body data if the content is embedded directly;
    attachment_id is set when the content must be fetched via attachments.get.
    """
    results: list[dict] = []
    filename = (payload.get("filename") or "").strip()
    body = payload.get("body") or {}
    if filename:
        results.append({
            "attachment_id": body.get("attachmentId", ""),
            "data_inline": body.get("data", ""),
            "filename": filename,
            "mime_type": payload.get("mimeType", ""),
            "size": body.get("size", 0),
        })
    for part in payload.get("parts", []):
        results.extend(_extract_attachments_metadata(part))
    return results


def download_attachment(credentials, message_id: str, attachment_id: str) -> bytes:
    """
    Download a Gmail attachment by its attachmentId and return raw decoded bytes.
    Requires gmail.readonly scope.
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    att = service.users().messages().attachments().get(
        userId="me",
        messageId=message_id,
        id=attachment_id,
    ).execute()
    data = att.get("data", "")
    return base64.urlsafe_b64decode(data + "==")


def extract_text_from_attachment(filename: str, mime_type: str, data_bytes: bytes) -> str | None:
    """
    Extract plain text from an attachment's raw bytes.
    Dispatches by MIME type and file extension.
    Returns extracted text (truncated to _ATTACHMENT_MAX_TEXT_BYTES) or None when
    the type is unsupported or parsing fails.  Never raises.
    """
    import io

    name_lower = (filename or "").lower()
    mime_lower = (mime_type or "").lower()

    try:
        # PDF
        if mime_lower == "application/pdf" or name_lower.endswith(".pdf"):
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(data_bytes))
            pages = [page.extract_text() or "" for page in reader.pages]
            text = "\n".join(pages).strip()
            return text[:_ATTACHMENT_MAX_TEXT_BYTES] or None

        # Word .docx
        if (
            mime_lower == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or name_lower.endswith(".docx")
        ):
            import docx as python_docx
            doc = python_docx.Document(io.BytesIO(data_bytes))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text[:_ATTACHMENT_MAX_TEXT_BYTES] or None

        # Excel .xlsx
        if (
            mime_lower == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            or name_lower.endswith(".xlsx")
        ):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data_bytes), read_only=True, data_only=True)
            rows: list[str] = []
            for sheet in wb.worksheets:
                rows.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        rows.append("\t".join(cells))
            text = "\n".join(rows).strip()
            return text[:_ATTACHMENT_MAX_TEXT_BYTES] or None

        # PowerPoint .pptx
        if (
            mime_lower == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or name_lower.endswith(".pptx")
        ):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data_bytes))
            slides: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                parts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                parts.append(t)
                if parts:
                    slides.append(f"[Slide {i}] " + " | ".join(parts))
            text = "\n".join(slides).strip()
            return text[:_ATTACHMENT_MAX_TEXT_BYTES] or None

        # Plain text variants
        if mime_lower.startswith("text/") or name_lower.endswith(
            (".txt", ".csv", ".json", ".md", ".html", ".xml", ".yaml", ".yml")
        ):
            text = data_bytes.decode("utf-8", errors="replace").strip()
            return text[:_ATTACHMENT_MAX_TEXT_BYTES] or None

        # Images — caller should note these separately; return None to signal unsupported
        return None

    except Exception as exc:
        logger.debug("extract_text_from_attachment failed for %s: %s", filename, exc)
        return None


def get_message_with_attachments(credentials, message_id: str) -> dict:
    """
    Fetch a single Gmail message including full body text and extracted attachment content.
    Returns the same shape as get_message_full plus an 'attachments' list where each
    entry contains {filename, mime_type, size} and either extracted_text or a note/error.
    Attachments are capped at _ATTACHMENT_MAX_COUNT per message.
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    m = service.users().messages().get(
        userId="me",
        id=message_id,
        format="full",
    ).execute()

    headers = {h["name"]: h["value"] for h in m.get("payload", {}).get("headers", [])}
    label_ids = m.get("labelIds", [])
    payload = m.get("payload", {})
    body = _extract_body(payload)

    msg: dict = {
        "id": m["id"],
        "threadId": m.get("threadId"),
        "from": headers.get("From", ""),
        "to": headers.get("To", ""),
        "subject": headers.get("Subject", ""),
        "date": headers.get("Date", ""),
        "snippet": m.get("snippet", ""),
        "body": body or m.get("snippet", ""),
        "is_read": "UNREAD" not in label_ids,
        "label_ids": label_ids,
    }

    att_meta = _extract_attachments_metadata(payload)
    attachments: list[dict] = []
    for meta in att_meta[:_ATTACHMENT_MAX_COUNT]:
        entry: dict = {
            "filename": meta["filename"],
            "mime_type": meta["mime_type"],
            "size": meta["size"],
        }
        mime_lower = (meta["mime_type"] or "").lower()
        name_lower = (meta["filename"] or "").lower()
        is_image = mime_lower.startswith("image/") or name_lower.endswith(
            (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg", ".tiff", ".ico")
        )
        if is_image:
            entry["note"] = "Image attachment — content not extracted"
        else:
            try:
                if meta.get("data_inline"):
                    raw_bytes = base64.urlsafe_b64decode(meta["data_inline"] + "==")
                elif meta.get("attachment_id"):
                    raw_bytes = download_attachment(credentials, message_id, meta["attachment_id"])
                else:
                    entry["note"] = "Attachment content unavailable"
                    attachments.append(entry)
                    continue
                text = extract_text_from_attachment(meta["filename"], meta["mime_type"], raw_bytes)
                if text:
                    entry["extracted_text"] = text
                else:
                    entry["note"] = "Unsupported attachment type — content not extracted"
            except Exception as exc:
                logger.warning(
                    "Attachment download/parse failed for %s in message %s: %s",
                    meta["filename"], message_id, exc,
                )
                entry["error"] = f"Could not extract content: {exc}"
        attachments.append(entry)

    msg["attachments"] = attachments
    return msg


def mark_message_unread(credentials, message_id: str) -> dict:
    """Add UNREAD label to a message (mark as unread)."""
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    return service.users().messages().modify(
        userId="me",
        id=message_id,
        body={"addLabelIds": ["UNREAD"]},
    ).execute()


def mark_message_read(credentials, message_id: str) -> dict:
    """Remove UNREAD label from a message (mark as read)."""
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    return service.users().messages().modify(
        userId="me",
        id=message_id,
        body={"removeLabelIds": ["UNREAD"]},
    ).execute()


def archive_message(credentials, message_id: str) -> dict:
    """Remove INBOX label from a message (archive it)."""
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    return service.users().messages().modify(
        userId="me",
        id=message_id,
        body={"removeLabelIds": ["INBOX"]},
    ).execute()


def get_user_email(credentials) -> str:
    """Return the authenticated Gmail user's email address."""
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    profile = service.users().getProfile(userId="me").execute()
    return profile.get("emailAddress", "")


# --------------------------------------------------------------------------
# Email Operations Agent helpers
# --------------------------------------------------------------------------
# Used by reply / reply_all / forward / draft-edit operations. All assume the
# OAuth credentials already grant the required scopes (gmail.compose for
# drafts, gmail.send for sending, gmail.readonly for thread/message reads,
# gmail.modify for trash).


def _parse_gmail_message_envelope(msg: dict) -> dict:
    """Extract subject/to/cc/bcc/from headers from a Gmail message dict (any format)."""
    headers: dict[str, str] = {}
    for h in msg.get("payload", {}).get("headers", []):
        name = (h.get("name") or "").lower()
        if name:
            headers[name] = h.get("value") or ""
    return {
        "subject": headers.get("subject", ""),
        "to": headers.get("to", ""),
        "cc": headers.get("cc", ""),
        "bcc": headers.get("bcc", ""),
        "from": headers.get("from", ""),
        "message_id": headers.get("message-id", ""),
        "in_reply_to": headers.get("in-reply-to", ""),
        "references": headers.get("references", ""),
    }


def _apply_threading_headers(mime, env: dict) -> None:
    """Set In-Reply-To / References on a reply so Gmail keeps it in the SAME
    conversation. Gmail threads by these RFC-2822 headers (plus a matching
    Re: subject) — supplying threadId alone is not enough and sends a new
    conversation. *env* is the parsed envelope of the message being replied to.
    """
    last_id = (env.get("message_id") or "").strip()
    if not last_id:
        return
    refs = (env.get("references") or "").strip()
    mime["In-Reply-To"] = last_id
    mime["References"] = f"{refs} {last_id}".strip() if refs else last_id


def _format_address_pairs(pairs: list[tuple[str, str]]) -> str:
    """Format [(name, email), ...] back into a comma-separated address string."""
    parts: list[str] = []
    for name, email in pairs:
        e = (email or "").strip()
        if not e:
            continue
        n = (name or "").strip()
        parts.append(f"{n} <{e}>" if n else e)
    return ", ".join(parts)


def _addr_pairs(addr_str: str) -> list[tuple[str, str]]:
    """Parse 'Name <email>, Name2 <email2>' into [(name, email), ...]; empties filtered."""
    if not addr_str:
        return []
    return [(n, e) for (n, e) in getaddresses([addr_str]) if e]


def _to_list_or_str(value) -> str:
    """Coerce a recipient field (str or list) into a comma-separated string."""
    if value is None:
        return ""
    if isinstance(value, list):
        return ", ".join(str(x).strip() for x in value if str(x).strip())
    return str(value).strip()


def _build_mime(body: str, html_body: str | None) -> MIMEText | MIMEMultipart:
    if html_body:
        m = MIMEMultipart("alternative")
        m.attach(MIMEText(body, "plain"))
        m.attach(MIMEText(html_body, "html"))
        return m
    return MIMEText(body, "plain")


def update_draft(
    credentials,
    draft_id: str,
    to: str | None = None,
    subject: str | None = None,
    body: str | None = None,
    cc: str | None = None,
    bcc: str | None = None,
    html_body: str | None = None,
) -> dict:
    """
    Replace fields on an existing Gmail draft. Any argument left as None preserves
    the draft's current value for that field. Returns the updated draft's id +
    final field values.
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    draft = service.users().drafts().get(userId="me", id=draft_id, format="full").execute()
    msg = draft.get("message", {})
    env = _parse_gmail_message_envelope(msg)

    new_subject = subject if subject is not None else env["subject"]
    new_to = to if to is not None else env["to"]
    new_cc = cc if cc is not None else env["cc"]
    new_bcc = bcc if bcc is not None else env["bcc"]

    existing_body = _extract_body(msg.get("payload", {})) or ""
    new_body = body if body is not None else existing_body

    mime = _build_mime(new_body, html_body)
    mime["to"] = new_to or ""
    mime["subject"] = new_subject or ""
    if new_cc:
        mime["cc"] = new_cc
    if new_bcc:
        mime["bcc"] = new_bcc

    raw = base64.urlsafe_b64encode(mime.as_bytes()).decode("utf-8")
    result = service.users().drafts().update(
        userId="me",
        id=draft_id,
        body={"message": {"raw": raw}},
    ).execute()

    logger.info("Gmail draft updated: id=%s subject=%s", result.get("id"), new_subject)
    return {
        "draft_id": result.get("id"),
        "subject": new_subject,
        "to": new_to,
        "cc": new_cc,
        "bcc": new_bcc,
    }


def add_recipients_to_draft(
    credentials,
    draft_id: str,
    to_add: list[str] | None = None,
    cc_add: list[str] | None = None,
    bcc_add: list[str] | None = None,
) -> dict:
    """Append recipients to an existing draft (deduplicated). Preserves existing entries."""
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    draft = service.users().drafts().get(userId="me", id=draft_id, format="full").execute()
    msg = draft.get("message", {})
    env = _parse_gmail_message_envelope(msg)
    body_text = _extract_body(msg.get("payload", {})) or ""

    existing_to = _addr_pairs(env["to"])
    existing_cc = _addr_pairs(env["cc"])
    existing_bcc = _addr_pairs(env["bcc"])

    seen = {e.lower() for _, e in (existing_to + existing_cc + existing_bcc)}

    def _append(pairs: list[tuple[str, str]], additions: list[str] | None) -> list[tuple[str, str]]:
        if not additions:
            return pairs
        out = list(pairs)
        for addr in additions:
            a = str(addr or "").strip()
            if not a:
                continue
            if a.lower() in seen:
                continue
            out.append(("", a))
            seen.add(a.lower())
        return out

    new_to = _append(existing_to, to_add)
    new_cc = _append(existing_cc, cc_add)
    new_bcc = _append(existing_bcc, bcc_add)

    mime = MIMEText(body_text, "plain")
    mime["to"] = _format_address_pairs(new_to)
    mime["subject"] = env["subject"]
    if new_cc:
        mime["cc"] = _format_address_pairs(new_cc)
    if new_bcc:
        mime["bcc"] = _format_address_pairs(new_bcc)

    raw = base64.urlsafe_b64encode(mime.as_bytes()).decode("utf-8")
    result = service.users().drafts().update(
        userId="me",
        id=draft_id,
        body={"message": {"raw": raw}},
    ).execute()

    return {
        "draft_id": result.get("id"),
        "to": _format_address_pairs(new_to),
        "cc": _format_address_pairs(new_cc),
        "bcc": _format_address_pairs(new_bcc),
        "subject": env["subject"],
        "added": {
            "to": list(to_add or []),
            "cc": list(cc_add or []),
            "bcc": list(bcc_add or []),
        },
    }


def remove_recipients_from_draft(
    credentials,
    draft_id: str,
    to_remove: list[str] | None = None,
    cc_remove: list[str] | None = None,
    bcc_remove: list[str] | None = None,
) -> dict:
    """Remove specific email addresses from a draft's recipient fields (case-insensitive)."""
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    draft = service.users().drafts().get(userId="me", id=draft_id, format="full").execute()
    msg = draft.get("message", {})
    env = _parse_gmail_message_envelope(msg)
    body_text = _extract_body(msg.get("payload", {})) or ""

    existing_to = _addr_pairs(env["to"])
    existing_cc = _addr_pairs(env["cc"])
    existing_bcc = _addr_pairs(env["bcc"])

    def _filter(pairs: list[tuple[str, str]], removals: list[str] | None) -> list[tuple[str, str]]:
        if not removals:
            return pairs
        rm = {str(r).strip().lower() for r in removals if str(r).strip()}
        return [(n, e) for (n, e) in pairs if e.lower() not in rm]

    new_to = _filter(existing_to, to_remove)
    new_cc = _filter(existing_cc, cc_remove)
    new_bcc = _filter(existing_bcc, bcc_remove)

    mime = MIMEText(body_text, "plain")
    mime["to"] = _format_address_pairs(new_to)
    mime["subject"] = env["subject"]
    if new_cc:
        mime["cc"] = _format_address_pairs(new_cc)
    if new_bcc:
        mime["bcc"] = _format_address_pairs(new_bcc)

    raw = base64.urlsafe_b64encode(mime.as_bytes()).decode("utf-8")
    result = service.users().drafts().update(
        userId="me",
        id=draft_id,
        body={"message": {"raw": raw}},
    ).execute()

    return {
        "draft_id": result.get("id"),
        "to": _format_address_pairs(new_to),
        "cc": _format_address_pairs(new_cc),
        "bcc": _format_address_pairs(new_bcc),
        "subject": env["subject"],
        "removed": {
            "to": list(to_remove or []),
            "cc": list(cc_remove or []),
            "bcc": list(bcc_remove or []),
        },
    }


def reply_to_thread(
    credentials,
    thread_id: str,
    body: str,
    html_body: str | None = None,
    cc: list[str] | str | None = None,
) -> dict:
    """
    Reply to the most recent message in a thread. Recipient defaults to the
    original sender; subject is preserved (with 'Re: ' prefix if missing).
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    thread = service.users().threads().get(
        userId="me",
        id=thread_id,
        format="metadata",
        metadataHeaders=["From", "To", "Cc", "Subject", "Message-ID", "References"],
    ).execute()
    messages = thread.get("messages", [])
    if not messages:
        raise ValueError(f"Thread {thread_id} has no messages")

    last = messages[-1]
    env = _parse_gmail_message_envelope(last)
    reply_to = env["from"]
    subject = env["subject"]
    if subject and not subject.lower().startswith("re:"):
        subject = "Re: " + subject

    mime = _build_mime(body, html_body)
    mime["to"] = reply_to
    mime["subject"] = subject
    cc_str = _to_list_or_str(cc)
    if cc_str:
        mime["cc"] = cc_str
    # Thread the reply: Gmail needs In-Reply-To + References matching the
    # thread's Message-IDs — threadId alone sends a NEW conversation.
    _apply_threading_headers(mime, env)

    raw = base64.urlsafe_b64encode(mime.as_bytes()).decode("utf-8")
    result = service.users().messages().send(
        userId="me",
        body={"raw": raw, "threadId": thread_id},
    ).execute()

    logger.info("Gmail reply sent: id=%s thread=%s to=%s", result.get("id"), thread_id, reply_to)
    return {
        "id": result.get("id"),
        "thread_id": thread_id,
        "to": reply_to,
        "subject": subject,
    }


def compute_reply_all_recipients(credentials, thread_id: str) -> dict:
    """
    Compute the To + Cc list a reply-all would use, WITHOUT sending anything.
    Mirrors the participant logic in reply_all_in_thread so the on-screen draft
    popup can display every recipient before the user confirms. The model cannot
    reliably aggregate/dedup participants across all thread messages itself, so
    the server is the source of truth.
    Returns {"to": str, "cc": list[str], "subject": str}.
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)

    profile = service.users().getProfile(userId="me").execute()
    my_email = (profile.get("emailAddress") or "").lower().strip()

    thread = service.users().threads().get(
        userId="me",
        id=thread_id,
        format="metadata",
        metadataHeaders=["From", "To", "Cc", "Subject", "Message-ID", "References"],
    ).execute()
    messages = thread.get("messages", [])
    if not messages:
        return {"to": "", "cc": [], "subject": ""}

    participants: dict[str, str] = {}
    for msg in messages:
        env = _parse_gmail_message_envelope(msg)
        for field in ("from", "to", "cc"):
            for name, email in _addr_pairs(env.get(field, "")):
                e_lc = email.lower()
                if not e_lc or e_lc == my_email:
                    continue
                if e_lc not in participants:
                    participants[e_lc] = f"{name} <{email}>" if name else email

    last_env = _parse_gmail_message_envelope(messages[-1])
    primary_to = ""
    primary_lc = ""
    for name, email in _addr_pairs(last_env.get("from", "")):
        if email.lower() != my_email:
            primary_to = f"{name} <{email}>" if name else email
            primary_lc = email.lower()
            break
    if not primary_to and participants:
        primary_lc, primary_to = next(iter(participants.items()))

    cc_list = [v for k, v in participants.items() if k != primary_lc]

    subject = last_env.get("subject", "")
    if subject and not subject.lower().startswith("re:"):
        subject = "Re: " + subject

    return {"to": primary_to, "cc": cc_list, "subject": subject}


def reply_all_in_thread(
    credentials,
    thread_id: str,
    body: str,
    html_body: str | None = None,
) -> dict:
    """
    Reply to every participant on a thread (To + Cc + From across all messages),
    excluding the authenticated user. The original sender becomes the primary To;
    everyone else becomes Cc.
    """
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)

    profile = service.users().getProfile(userId="me").execute()
    my_email = (profile.get("emailAddress") or "").lower().strip()

    thread = service.users().threads().get(
        userId="me",
        id=thread_id,
        format="metadata",
        metadataHeaders=["From", "To", "Cc", "Subject", "Message-ID", "References"],
    ).execute()
    messages = thread.get("messages", [])
    if not messages:
        raise ValueError(f"Thread {thread_id} has no messages")

    # Collect every distinct participant, preserving the formatted "Name <email>" form.
    participants: dict[str, str] = {}
    for msg in messages:
        env = _parse_gmail_message_envelope(msg)
        for field in ("from", "to", "cc"):
            for name, email in _addr_pairs(env.get(field, "")):
                e_lc = email.lower()
                if not e_lc or e_lc == my_email:
                    continue
                if e_lc not in participants:
                    participants[e_lc] = f"{name} <{email}>" if name else email

    # Primary To = the From of the last inbound message (excluding self).
    last_env = _parse_gmail_message_envelope(messages[-1])
    primary_to = ""
    primary_lc = ""
    for name, email in _addr_pairs(last_env.get("from", "")):
        if email.lower() != my_email:
            primary_to = f"{name} <{email}>" if name else email
            primary_lc = email.lower()
            break

    if not primary_to and participants:
        # Fallback: first remaining participant.
        primary_lc, primary_to = next(iter(participants.items()))

    cc_list = [v for k, v in participants.items() if k != primary_lc]

    subject = last_env.get("subject", "")
    if subject and not subject.lower().startswith("re:"):
        subject = "Re: " + subject

    mime = _build_mime(body, html_body)
    mime["to"] = primary_to
    mime["subject"] = subject
    if cc_list:
        mime["cc"] = ", ".join(cc_list)
    # Keep the reply-all in the SAME conversation: set In-Reply-To + References
    # from the last message. Without these Gmail starts a new thread even when
    # threadId is supplied.
    _apply_threading_headers(mime, last_env)

    raw = base64.urlsafe_b64encode(mime.as_bytes()).decode("utf-8")
    result = service.users().messages().send(
        userId="me",
        body={"raw": raw, "threadId": thread_id},
    ).execute()

    logger.info(
        "Gmail reply-all sent: id=%s thread=%s to=%s cc_count=%s",
        result.get("id"),
        thread_id,
        primary_to,
        len(cc_list),
    )
    return {
        "id": result.get("id"),
        "thread_id": thread_id,
        "to": primary_to,
        "cc": ", ".join(cc_list) if cc_list else "",
        "subject": subject,
        "recipient_count": (1 if primary_to else 0) + len(cc_list),
        "all_recipients": ([primary_to] if primary_to else []) + cc_list,
    }


def forward_message(
    credentials,
    message_id: str,
    to,
    body: str | None = None,
    html_body: str | None = None,
) -> dict:
    """
    Forward an existing inbox message to new recipients. Builds the standard
    'Forwarded message' envelope inline; attachments are NOT carried (text only
    for now).
    """
    original = get_message_full(credentials, message_id)
    orig_subject = original.get("subject") or ""
    new_subject = orig_subject if orig_subject.lower().startswith("fwd:") else (
        f"Fwd: {orig_subject}" if orig_subject else "Fwd:"
    )

    orig_body = original.get("body") or original.get("snippet") or ""
    user_prefix = (body or "").strip()
    forwarded = (
        (f"{user_prefix}\n\n" if user_prefix else "")
        + "---------- Forwarded message ----------\n"
        + f"From: {original.get('from', '')}\n"
        + f"Date: {original.get('date', '')}\n"
        + f"Subject: {orig_subject}\n"
        + f"To: {original.get('to', '')}\n\n"
        + f"{orig_body}\n"
    )

    to_str = _to_list_or_str(to)
    if not to_str:
        raise ValueError("Forward requires at least one recipient.")

    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    mime = _build_mime(forwarded, html_body)
    mime["to"] = to_str
    mime["subject"] = new_subject

    raw = base64.urlsafe_b64encode(mime.as_bytes()).decode("utf-8")
    result = service.users().messages().send(
        userId="me",
        body={"raw": raw},
    ).execute()

    logger.info(
        "Gmail forward sent: id=%s forwarded_message_id=%s to=%s",
        result.get("id"),
        message_id,
        to_str,
    )
    return {
        "id": result.get("id"),
        "to": to_str,
        "subject": new_subject,
        "forwarded_message_id": message_id,
    }


def trash_message(credentials, message_id: str) -> dict:
    """Move a message to Gmail Trash (recoverable for 30 days)."""
    service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
    result = service.users().messages().trash(userId="me", id=message_id).execute()
    logger.info("Gmail message trashed: id=%s", message_id)
    return {"id": message_id, "status": "trashed", "label_ids": result.get("labelIds", [])}
